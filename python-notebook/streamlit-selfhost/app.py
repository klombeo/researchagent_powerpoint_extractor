import streamlit as st
import zipfile
import os
import shutil
import tempfile
import io
from pptx import Presentation
import openpyxl
from openpyxl.styles import numbers
from PIL import Image
import xml.etree.ElementTree as ET

# ----------------- Helper Functions -----------------

def map_embedded_files_to_slides(pptx_file):
    rels_map = {}
    with zipfile.ZipFile(pptx_file, 'r') as z:
        slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
        for slide_file in slide_files:
            slide_index = int(slide_file.split("slide")[-1].split(".xml")[0]) - 1
            rel_path = f"ppt/slides/_rels/slide{slide_index+1}.xml.rels"
            if rel_path in z.namelist():
                rel_data = z.read(rel_path)
                tree = ET.fromstring(rel_data)
                for rel in tree.findall(".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                    target = rel.attrib.get("Target")
                    if target and target.startswith("../embeddings/"):
                        embedded_file = "ppt/embeddings/" + os.path.basename(target)
                        rels_map[embedded_file] = slide_index
    return rels_map

def extract_excel_from_pptx(pptx_file, output_dir, base_filename):
    with zipfile.ZipFile(pptx_file, 'r') as z:
        embedded_files = [f for f in z.namelist() if f.startswith('ppt/embeddings/') and (f.endswith('.xlsx') or f.endswith('.xls'))]
        total_extracted = 0

        prs = Presentation(pptx_file)
        slide_titles = []
        for slide in prs.slides:
            title = ""
            subtitle = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not title:
                        title = text
                    elif not subtitle:
                        subtitle = text
            slide_titles.append(f"{title} - {subtitle}" if subtitle else title)

        embedded_map = map_embedded_files_to_slides(pptx_file)

        for embedded_file in embedded_files:
            file_data = z.read(embedded_file)
            ext = ".xlsx" if embedded_file.endswith(".xlsx") else ".xls"
            index = embedded_files.index(embedded_file)
            out_filename = f"{base_filename}_embedded_excel_{index+1}{ext}"
            out_path = os.path.join(output_dir, out_filename)

            with open(out_path, 'wb') as f:
                f.write(file_data)
                total_extracted += 1

            if ext == ".xlsx":
                try:
                    wb = openpyxl.load_workbook(out_path)
                    for sheet in wb.worksheets:
                        max_col = sheet.max_column
                        sheet.insert_cols(max_col + 1)
                        sheet.cell(row=1, column=max_col + 1).value = "table_title"
                        slide_idx = embedded_map.get(embedded_file, index)
                        title_val = slide_titles[slide_idx] if slide_idx < len(slide_titles) else ""
                        for row in range(2, sheet.max_row + 1):
                            sheet.cell(row=row, column=max_col + 1).value = title_val
                    wb.save(out_path)
                except Exception as e:
                    print("Error adding table_title column:", e)

    return total_extracted

# (rest of the code remains unchanged)

# NOTE: Make sure the rest of your application code (cleaning, UI, image extraction) follows below

def clean_excel_files_in_folder(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".xlsx"):
            filepath = os.path.join(folder_path, filename)
            wb = openpyxl.load_workbook(filepath)
            for sheet in wb.worksheets:
                max_column = sheet.max_column
                headers = [cell.value for cell in sheet[1]]
                for row in sheet.iter_rows(min_row=2):
                    for idx, cell in enumerate(row):
                        if idx < len(headers) and headers[idx] is not None:
                            if cell.value is None:
                                cell.value = 0
                                cell.number_format = '0'
                            elif isinstance(cell.value, str):
                                val = cell.value.strip().lower()
                                if val in ['n/a', 'no data', 'undefined', '-', '']:
                                    cell.value = 0
                                    cell.number_format = '0'
            wb.save(filepath)

def extract_images_from_pptx(pptx_path, output_dir, base_filename):
    prs = Presentation(pptx_path)
    image_count = 0
    safe_base = base_filename.replace(" ", "_")

    for slide_num, slide in enumerate(prs.slides):
        for shape_num, shape in enumerate(slide.shapes):
            if shape.shape_type == 13 and hasattr(shape, "image"):
                image = shape.image
                ext = image.ext
                image_data = image.blob
                filename = f"{safe_base}_image_{image_count+1}.{ext}"
                image_path = os.path.join(output_dir, filename)

                with open(image_path, "wb") as f:
                    f.write(image_data)
                    image_count += 1

    return image_count

# ----------------- Streamlit App -----------------

st.set_page_config(page_title="PowerPoint Extractor", layout="centered")

st.markdown("""
    <style>
    button[data-testid="base-button"] {
        background-color: #4CAF50;
        color: white;
        padding: 0.4rem 1rem;
        font-size: 14px;
        border-radius: 10px;
        border: none;
        transition: 0.1s ease-in-out;
        width: auto;
    }
    button[data-testid="base-button"]:hover {
        background-color: #388e3c;
        transform: scale(1.02);
    }
    button[data-testid="base-button"]:active {
        transform: scale(0.95);
    }
    </style>
""", unsafe_allow_html=True)

st.title("📤 PowerPoint Excel/Image Extractor")

uploaded_files = st.file_uploader("Upload one or more .pptx files", type=["pptx"], accept_multiple_files=True)

if uploaded_files:
    extract_option = st.radio("What do you want to extract?", ("Excel", "Images"))

    if st.button("🔍 Extract Selected Content"):
        with tempfile.TemporaryDirectory() as tmpdir:
            combined_dir = os.path.join(tmpdir, "extracted")
            os.makedirs(combined_dir, exist_ok=True)
            total_all = 0
            output_zip_path = os.path.join(tmpdir, "output.zip")
            image_dir = ""

            for uploaded_file in uploaded_files:
                pptx_path = os.path.join(tmpdir, uploaded_file.name)
                with open(pptx_path, "wb") as f:
                    f.write(uploaded_file.read())

                base_filename = os.path.splitext(uploaded_file.name)[0].replace(" ", "_")

                if extract_option == "Excel":
                    excel_dir = os.path.join(combined_dir, "excel")
                    os.makedirs(excel_dir, exist_ok=True)

                    total = extract_excel_from_pptx(pptx_path, excel_dir, base_filename)
                    clean_excel_files_in_folder(excel_dir)
                    total_all += total

                elif extract_option == "Images":
                    image_dir = os.path.join(combined_dir, "images")
                    os.makedirs(image_dir, exist_ok=True)

                    total = extract_images_from_pptx(pptx_path, image_dir, base_filename)
                    total_all += total

            if total_all == 0:
                st.warning(f"No {extract_option.lower()} files found in the uploaded PowerPoint files.")
            else:
                st.success(f"Successfully extracted {total_all} {extract_option.lower()} file(s).")

                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "w") as zipf:
                    for root, _, files in os.walk(combined_dir):
                        for file in files:
                            full_path = os.path.join(root, file)
                            arcname = os.path.relpath(full_path, combined_dir)
                            zipf.write(full_path, arcname=arcname)

                label = f"📥 Download {extract_option} Files (.zip)"
                st.download_button(label, data=zip_buffer.getvalue(), file_name=f"pptx_{extract_option.lower()}_output.zip", mime="application/zip")

                if extract_option == "Images":
                    st.subheader("🖼️ Preview of Extracted Images")
                    if image_dir:
                        for file in os.listdir(image_dir):
                            img_path = os.path.join(image_dir, file)
                            image = Image.open(img_path)
                            st.image(image, caption=file, use_container_width=True)

                        with open(output_zip_path, "wb") as out_f:
                            out_f.write(zip_buffer.getvalue())

                        with open(output_zip_path, "rb") as f:
                            st.download_button("📸 Download Extracted Images (.zip)", f.read(), file_name="extracted_images.zip", mime="application/zip")
